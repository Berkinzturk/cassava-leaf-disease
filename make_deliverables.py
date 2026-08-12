"""
Builds the three submission files from whatever the notebook produced, then zips them.

    pip install reportlab python-pptx
    python make_deliverables.py

Reads:  results.json, figures/*.png, cassava_capstone.ipynb
Writes: report.pdf, slides.pptx, Capstone_Cassava_<lastname>.zip

If results.json is missing, it still builds both documents with "[run notebook]" in
every numeric slot, so you can check the layout before training finishes.
"""

import glob
import json
import os
import zipfile
from pathlib import Path

# ------------------------------------------------------------------ config
# EDIT THIS after you create the repo. It goes on the title page and slide 1.
GITHUB_URL = "GITHUB_LINK_HERE"

STUDENT = "Berkin Ozturk"
COURSE = "Deep Learning with PyTorch"
TITLE = "Cassava Leaf Disease Classification"
SUBTITLE = "A CNN classifier with a language-model advice layer"

HERE = Path(__file__).resolve().parent
FIG = HERE / "figures"
OUT_PDF = HERE / "report.pdf"
OUT_PPTX = HERE / "slides.pptx"
OUT_ZIP = HERE / "Capstone_Cassava_Ozturk.zip"

# Forest and moss. Chosen because the subject is leaves.
FOREST = "2C5F2D"
MOSS = "97BC62"
CREAM = "F5F5F5"
INK = "1A1A1A"
GREY = "6B6B6B"

MISSING = "[run notebook]"

CLASS_LABELS = {
    "cbb": "Bacterial Blight",
    "cbsd": "Brown Streak",
    "cgm": "Green Mite",
    "cmd": "Mosaic Disease",
    "healthy": "Healthy",
}

# ------------------------------------------------------------------ results

R = {}
if (HERE / "results.json").exists():
    R = json.loads((HERE / "results.json").read_text())
    print("loaded results.json")
else:
    print("no results.json found, building with placeholders")


def pct(key, digits=1):
    v = R.get(key)
    return MISSING if v is None else f"{100 * float(v):.{digits}f}%"


def num(key, digits=3):
    v = R.get(key)
    return MISSING if v is None else f"{float(v):.{digits}f}"


def gap(unit=""):
    a, b = R.get("resnet18_test_acc"), R.get("smallcnn_test_acc")
    if a is None or b is None:
        return MISSING
    return f"{100 * (float(a) - float(b)):+.1f}{unit}"


def classes():
    return R.get("classes") or ["cbb", "cbsd", "cgm", "cmd", "healthy"]


def per_class_rows():
    pc = R.get("per_class") or {}
    rows = [["Class", "Precision", "Recall", "F1", "Support"]]
    for c in classes():
        d = pc.get(c)
        if d:
            rows.append([CLASS_LABELS.get(c, c), f"{d['precision']:.3f}",
                         f"{d['recall']:.3f}", f"{d['f1-score']:.3f}",
                         str(int(d['support']))])
        else:
            rows.append([CLASS_LABELS.get(c, c), MISSING, MISSING, MISSING, MISSING])
    return rows


def sweep_rows():
    rows = [["Learning rate", "Batch size", "Weight decay", "Val macro F1"]]
    sw = R.get("sweep")
    if not sw:
        for lr, bs, wd in [("1e-3", "32", "1e-4"), ("3e-4", "32", "1e-4"),
                           ("1e-4", "32", "1e-4"), ("3e-4", "64", "1e-4"),
                           ("3e-4", "32", "1e-2")]:
            rows.append([lr, bs, wd, MISSING])
        return rows
    for s in sorted(sw, key=lambda d: -d["val_macro_f1"]):
        rows.append([f"{s['lr']:g}", str(int(s["batch_size"])),
                     f"{s['weight_decay']:g}", f"{s['val_macro_f1']:.4f}"])
    return rows


def best_cfg_text():
    b = R.get("best_config")
    if not b:
        return MISSING
    return (f"learning rate {float(b['lr']):g}, batch size {int(b['batch_size'])}, "
            f"weight decay {float(b['weight_decay']):g}")


def fig(name):
    p = FIG / name
    return str(p) if p.exists() else None


# ================================================================== PDF

def build_pdf():
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_JUSTIFY
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (Image, PageBreak, Paragraph, SimpleDocTemplate,
                                    Spacer, Table, TableStyle)

    # Try for real fonts. Colab ships DejaVu with matplotlib; fall back to the
    # built-ins if this machine does not have them.
    body_font, head_font, bold_font = "Times-Roman", "Helvetica", "Times-Bold"
    cand = glob.glob("/usr/share/fonts/**/DejaVuSerif.ttf", recursive=True)
    cand += glob.glob("/usr/**/matplotlib/mpl-data/fonts/ttf/DejaVuSerif.ttf", recursive=True)
    try:
        import matplotlib
        mpl = Path(matplotlib.__file__).parent / "mpl-data" / "fonts" / "ttf"
        cand.append(str(mpl / "DejaVuSerif.ttf"))
    except Exception:
        pass
    for serif in cand:
        s = Path(serif)
        if s.exists() and (s.parent / "DejaVuSerif-Bold.ttf").exists() \
                and (s.parent / "DejaVuSans.ttf").exists():
            pdfmetrics.registerFont(TTFont("BodySerif", str(s)))
            pdfmetrics.registerFont(TTFont("BodySerif-Bold", str(s.parent / "DejaVuSerif-Bold.ttf")))
            pdfmetrics.registerFont(TTFont("HeadSans", str(s.parent / "DejaVuSans-Bold.ttf")))
            body_font, bold_font, head_font = "BodySerif", "BodySerif-Bold", "HeadSans"
            break
    print("pdf fonts:", body_font, "/", head_font)

    green = colors.HexColor("#" + FOREST)
    ink = colors.HexColor("#" + INK)
    grey = colors.HexColor("#" + GREY)

    S = {
        "title": ParagraphStyle("t", fontName=head_font, fontSize=26, leading=31,
                                textColor=green, spaceAfter=6),
        "sub": ParagraphStyle("s", fontName=body_font, fontSize=13, leading=17,
                              textColor=grey, spaceAfter=26),
        "h1": ParagraphStyle("h1", fontName=head_font, fontSize=15, leading=19,
                             textColor=green, spaceBefore=16, spaceAfter=7),
        "h2": ParagraphStyle("h2", fontName=head_font, fontSize=11.5, leading=15,
                             textColor=ink, spaceBefore=11, spaceAfter=4),
        "p": ParagraphStyle("p", fontName=body_font, fontSize=10, leading=15.5,
                            textColor=ink, alignment=TA_JUSTIFY, spaceAfter=8),
        "meta": ParagraphStyle("m", fontName=body_font, fontSize=10.5, leading=17,
                               textColor=ink),
        "cap": ParagraphStyle("c", fontName=body_font, fontSize=8.5, leading=11,
                              textColor=grey, spaceBefore=3, spaceAfter=12),
        "code": ParagraphStyle("cd", fontName="Courier", fontSize=8.5, leading=12,
                               textColor=ink, spaceAfter=8),
    }

    def P(t, s="p"):
        return Paragraph(t, S[s])

    def table(rows, widths=None, small=False, header=True):
        t = Table(rows, colWidths=widths, hAlign="LEFT", repeatRows=1 if header else 0)
        style = [
            ("FONTNAME", (0, 0), (-1, 0), head_font),
            ("FONTNAME", (0, 1), (-1, -1), body_font),
            ("FONTSIZE", (0, 0), (-1, -1), 8 if small else 9),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("BACKGROUND", (0, 0), (-1, 0), green),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("LEFTPADDING", (0, 0), (-1, -1), 7),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F1F4EE")]),
            ("LINEBELOW", (0, -1), (-1, -1), 0.6, green),
        ]
        if not header:
            # drop the header-row rules, or row 0 renders white on white
            style = [x for x in style if x[0] not in ("BACKGROUND", "TEXTCOLOR")]
            style += [("FONTNAME", (0, 0), (-1, -1), body_font),
                      ("TEXTCOLOR", (0, 0), (-1, -1), ink),
                      ("FONTNAME", (0, 0), (0, -1), head_font),
                      ("TEXTCOLOR", (0, 0), (0, -1), green),
                      ("ROWBACKGROUNDS", (0, 0), (-1, -1),
                       [colors.white, colors.HexColor("#F1F4EE")])]
        t.setStyle(TableStyle(style))
        return t

    def picture(name, caption, width=6.2):
        p = fig(name)
        out = []
        if p:
            from PIL import Image as PILImage
            w, h = PILImage.open(p).size
            iw = width * inch
            out.append(Image(p, width=iw, height=iw * h / w))
        else:
            out.append(P(f"<i>[figure {name} appears here after the notebook run]</i>", "cap"))
        out.append(P(caption, "cap"))
        return out

    doc = SimpleDocTemplate(str(OUT_PDF), pagesize=LETTER,
                            leftMargin=0.95 * inch, rightMargin=0.95 * inch,
                            topMargin=0.9 * inch, bottomMargin=0.85 * inch,
                            title=TITLE, author=STUDENT)

    def footer(canvas, d):
        canvas.saveState()
        canvas.setFont(body_font, 8)
        canvas.setFillColor(grey)
        canvas.drawString(0.95 * inch, 0.5 * inch, f"{STUDENT} | {TITLE}")
        canvas.drawRightString(LETTER[0] - 0.95 * inch, 0.5 * inch, str(d.page))
        canvas.restoreState()

    s = []

    # ---- title page
    s += [Spacer(1, 1.5 * inch), P(TITLE, "title"), P(SUBTITLE, "sub")]
    s += [table([
        ["Student", STUDENT],
        ["Course", COURSE],
        ["Assignment", "Capstone project"],
        ["Dataset", "iCassava 2019 (Mwebaze et al.)"],
        ["Task", "5-class image classification"],
        ["Code", GITHUB_URL],
    ], widths=[1.4 * inch, 4.6 * inch], header=False)]
    s += [Spacer(1, 0.35 * inch)]
    if GITHUB_URL == "GITHUB_LINK_HERE":
        s += [P("<b>Reminder to self: replace GITHUB_LINK_HERE in make_deliverables.py "
                "and rebuild before submitting.</b>", "cap")]
    s += [PageBreak()]

    # ---- 1
    s += [P("1. Problem and motivation", "h1")]
    s += [P(
        "Cassava is a staple crop across sub-Saharan Africa, and four diseases account for "
        "most of the yield lost to plant pathology: cassava mosaic disease, brown streak "
        "disease, bacterial blight, and green mite damage. Catching them early changes the "
        "outcome, but expert diagnosis is scarce in the places where the crop matters most, "
        "and the visual signs overlap enough that a non-expert struggles to tell them apart.")]
    s += [P(
        "I picked this problem for two reasons. First, the images are real. They were taken by "
        "farmers on phones in the field, so they have bad lighting, odd angles, and cluttered "
        "backgrounds. That makes it a harder problem than the lab-photo plant datasets where a "
        "leaf sits flat on white paper and a model can hit 99% without learning much. Second, "
        "a class label on its own does not help a farmer. Turning that label into a written "
        "recommendation is where the language-model half of this project comes in.")]

    s += [P("2. Dataset", "h1")]
    s += [P(
        "The dataset is iCassava 2019, collected by the Artificial Intelligence lab at Makerere "
        "University together with the National Crops Resources Research Institute in Uganda. "
        "Plant pathologists labelled every image. It contains 9,430 labelled photos across five "
        "classes and ships with a fixed train, validation, and test split of 5,656 / 1,889 / "
        "1,885 images. I used that split as given rather than reshuffling, so the results here "
        "are comparable to other work on the same data.")]
    s += [P(
        "This dataset was not used in any lecture or homework for this course, and it is not one "
        "of the built-in torchvision datasets. It downloads from a public URL, so the notebook "
        "needs no Kaggle account or credentials to reproduce the run.")]
    s += picture("class_distribution.png",
                 "Figure 1. Training images per class. The imbalance is the central difficulty "
                 "in this dataset.")
    s += [P(
        "The class distribution is heavily skewed. Mosaic disease has roughly ten times the "
        "images of the healthy class. This matters more than it first appears. A model that "
        "predicts mosaic disease for every image scores respectably on accuracy while being "
        "useless, so accuracy is a misleading headline number here. Two decisions follow from "
        "that: the loss function is weighted by inverse class frequency, and every model "
        "selection decision in this project is made on macro F1 rather than accuracy.")]

    s += [P("3. Preprocessing and augmentation", "h1")]
    s += [P(
        "Images are resized to 224 by 224 and normalized with ImageNet channel statistics, "
        "since the transfer model was pretrained on that distribution. The augmentation is "
        "deliberately plain and matched to how the photos vary in reality: random resized crop "
        "down to 60% of the frame, horizontal and vertical flips, rotation up to 20 degrees, "
        "colour jitter on brightness, contrast, saturation and hue, and random erasing. A leaf "
        "has no canonical orientation, so vertical flips are safe here in a way they would not "
        "be for, say, digit recognition. The validation and test transforms do none of this: "
        "resize to 256, centre crop to 224, normalize.")]
    s += picture("augmented_batch.png",
                 "Figure 2. One augmented training batch. The variation in crop, angle, and "
                 "colour is what the model sees during training.")

    s += [PageBreak()]

    # ---- 4
    s += [P("4. Model architectures", "h1")]
    s += [P(
        "I trained two models so there would be something to compare against.")]
    s += [P("SmallCNN, built from scratch", "h2")]
    s += [P(
        "Four convolutional blocks, each with two 3 by 3 convolutions, batch normalization, "
        "ReLU, and max pooling, widening from 32 to 256 channels. Global average pooling, "
        "dropout at 0.4, then a linear layer to five outputs. Around 1.2 million parameters. "
        "This is the honest baseline: it shows what the dataset gives you with no outside help.")]
    s += [P("ResNet18, fine-tuned", "h2")]
    s += [P(
        "Standard ResNet18 with ImageNet weights and the final fully connected layer replaced "
        "by a fresh five-way head. The whole network is trainable, not just the head. Around "
        "11.2 million parameters. With fewer than 6,000 training images, learning good "
        "low-level filters from random initialization is not realistic, so I expected this to "
        "win by a wide margin. The size of that gap is the more interesting result than either "
        "model's absolute score.")]
    s += [P(
        "Regularization comes from four places: dropout in the small model, weight decay in "
        "AdamW, label smoothing at 0.05, and the augmentation described above. Early stopping "
        "watches validation macro F1 with a patience of five epochs and restores the weights "
        "from the best epoch rather than the last one.")]

    s += [P("5. Training setup", "h1")]
    s += [table([
        ["Setting", "Value", "Reason"],
        ["Optimizer", "AdamW", "Decoupled weight decay behaves better than Adam here"],
        ["Schedule", "OneCycleLR", "Warms up then anneals, converges in fewer epochs"],
        ["Loss", "Cross entropy, class-weighted", "Inverse frequency, so rare classes count"],
        ["Label smoothing", "0.05", "Stops the model becoming overconfident"],
        ["Precision", "Mixed (AMP)", "Roughly doubles throughput on a T4"],
        ["Early stop metric", "Validation macro F1", "Accuracy hides rare-class collapse"],
        ["Seed", "42", "Set on Python, NumPy, and Torch"],
    ], widths=[1.5 * inch, 1.7 * inch, 3.1 * inch])]
    s += [Spacer(1, 10)]
    s += picture("training_curves.png",
                 "Figure 3. Loss, accuracy, and macro F1 for both models. Dashed lines are "
                 "training, solid lines are validation.")

    s += [PageBreak()]

    # ---- 6
    s += [P("6. Hyperparameter tuning", "h1")]
    s += [P(
        "I ran a five-point sweep on the ResNet over learning rate, batch size, and weight "
        "decay. Each configuration trained for six epochs, which is enough to rank them without "
        "using up a free Colab session. The best configuration was then retrained properly for "
        "up to 25 epochs with early stopping.")]
    s += [table(sweep_rows(), widths=[1.5 * inch, 1.4 * inch, 1.5 * inch, 1.6 * inch])]
    s += [Spacer(1, 8)]
    s += [P(f"Best configuration: {best_cfg_text()}.")]
    s += [P(
        "The sweep is short and single-seed, so small differences between neighbouring rows are "
        "not meaningful. It tells you which region of the space to be in, not which exact "
        "setting is optimal.")]

    s += [P("7. Results", "h1")]
    s += [table([
        ["Model", "Test accuracy", "Macro precision", "Macro recall", "Macro F1"],
        ["SmallCNN (scratch)", pct("smallcnn_test_acc"),
         num("smallcnn_test_macro_precision"), num("smallcnn_test_macro_recall"),
         num("smallcnn_test_macro_f1")],
        ["ResNet18 (fine-tuned)", pct("resnet18_test_acc"),
         num("resnet18_test_macro_precision"), num("resnet18_test_macro_recall"),
         num("resnet18_test_macro_f1")],
    ], widths=[1.75 * inch, 1.15 * inch, 1.15 * inch, 1.05 * inch, 0.9 * inch], small=True)]
    s += [Spacer(1, 8)]
    s += [P(f"Transfer learning moved test accuracy by {gap(' points')} over the from-scratch CNN. "
            "That is the headline finding, and it is the expected one: on a dataset this size, "
            "pretrained features are worth more than any architecture tweak I could make.")]
    s += [P("Per-class performance, ResNet18 on the test set", "h2")]
    s += [table(per_class_rows(),
                widths=[1.6 * inch, 1.1 * inch, 1.0 * inch, 1.0 * inch, 1.0 * inch])]
    s += [Spacer(1, 10)]
    s += picture("confusion_test_resnet18.png",
                 "Figure 4. Confusion matrix, ResNet18 on the held-out test set.", width=4.6)

    s += [PageBreak()]

    s += [P("8. Error analysis", "h1")]
    s += [P(
        "The mistakes worth looking at are the confident ones. A model that is wrong at 40% "
        "confidence is at least signalling uncertainty; one that is wrong at 95% is not.")]
    s += picture("mistakes.png",
                 "Figure 5. The twelve misclassified validation images the model was most "
                 "confident about.")
    s += [P(
        "Write two or three sentences here after looking at Figure 5 in your own run. Things "
        "worth checking: whether the confused pairs are the diseases that genuinely look alike, "
        "whether the healthy class is being lost to the majority classes, and whether any of "
        "the errors are images a human would also get wrong or would call mislabelled.")]

    s += [P("9. Passing the output to a language model", "h1")]
    s += [P(
        "A five-way softmax is not something a farmer can act on. The second half of this "
        "project sends the classifier's top three predictions, with their probabilities, to "
        "Claude through the Anthropic API. The response is a short set of next steps written "
        "for a non-specialist.")]
    s += [P(
        "The probabilities go into the prompt on purpose. When the classifier is genuinely "
        "unsure, the language model can say so and recommend a human check, instead of "
        "converting a coin flip into confident-sounding advice. That is the one design decision "
        "in this part that does real work.")]
    s += [P("The call itself is five lines:")]
    s += [P("msg = client.messages.create(<br/>"
            "&nbsp;&nbsp;&nbsp;&nbsp;model=LLM_MODEL,<br/>"
            "&nbsp;&nbsp;&nbsp;&nbsp;max_tokens=400,<br/>"
            "&nbsp;&nbsp;&nbsp;&nbsp;messages=[{&quot;role&quot;: &quot;user&quot;, "
            "&quot;content&quot;: prompt}])<br/>"
            "return msg.content[0].text", "code")]

    s += [P("10. Interface", "h1")]
    s += [P(
        "The whole thing is wrapped in a Gradio interface. Upload a leaf photo, and the page "
        "returns the class probabilities as a bar display alongside the written recommendation. "
        "It runs from the notebook with a public share link, and there is also a standalone "
        "app.py in the repository that loads the saved checkpoint.")]
    s += picture("ui_screenshot.png",
                 "Figure 6. The Gradio interface. Save a screenshot as "
                 "figures/ui_screenshot.png and rebuild to include it.", width=5.6)

    s += [PageBreak()]

    s += [P("11. What this does not prove", "h1")]
    s += [P(
        "The test set comes from the same collection effort as the training set. Same region, "
        "same season, same set of phones, same field conditions. So a good test score here is "
        "evidence that the model learned the visual patterns in this collection, and it is not "
        "evidence that the model works on a different farm, a different camera, or a different "
        "time of year. Nothing in this project measures that, and the standard way to measure "
        "it would be a field trial with newly collected images.")]
    s += [P(
        "The language model output is unverified. It reads well and it is plausible, but no "
        "agronomist reviewed a single response, and I have no way to score whether the advice "
        "is agriculturally correct. Worse, the system produces equally fluent advice when the "
        "classifier is wrong. A wrong label is a small failure; a wrong label wrapped in "
        "confident treatment instructions is a larger one, because it is more likely to be "
        "acted on. Anyone deploying this would need expert review of the generated text and a "
        "confidence threshold below which the system refuses to advise at all.")]
    s += [P(
        "The hyperparameter sweep is also thin. Five configurations, six epochs each, one seed. "
        "It narrows the region but it does not establish that the chosen setting is best, and "
        "run-to-run variance was never measured.")]

    s += [P("12. Conclusion and future work", "h1")]
    s += [P(
        "The project does what it set out to do: a fine-tuned ResNet18 separates five cassava "
        "leaf conditions on real field photography, clearly outperforms an equivalent network "
        "trained from scratch, and the output reaches the user as advice rather than as a class "
        "index. The imbalance handling and the choice of macro F1 as the selection metric are "
        "the two decisions that most affected the outcome.")]
    s += [P("Things I would do next, roughly in order of expected value:")]
    for item in [
        "Test on images from a different source to get an honest read on generalization.",
        "Add a confidence threshold so low-certainty cases route to a human instead of the "
        "language model.",
        "Try a stronger backbone (EfficientNet-B0 or ConvNeXt-Tiny) now that the pipeline works.",
        "Add Grad-CAM so the interface can show which part of the leaf drove the prediction, "
        "which also makes it easier to catch the model latching onto background artifacts.",
        "Repeat the best configuration across several seeds to see how much of the gap between "
        "sweep rows is noise.",
    ]:
        s += [P(f"&bull;&nbsp;&nbsp;{item}")]

    s += [P("Reference", "h1")]
    s += [P("Mwebaze, E., Gebru, T., Frome, A., Nsumba, S., and Tusubira, J. (2019). "
            "<i>iCassava 2019 Fine-Grained Visual Categorization Challenge.</i> "
            "arXiv:1908.02900.")]
    s += [P(f"Code and full notebook: {GITHUB_URL}")]

    doc.build(s, onFirstPage=lambda c, d: None, onLaterPages=footer)
    print("wrote", OUT_PDF.name)


# ================================================================== PPTX

def build_pptx():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    def rgb(h):
        return RGBColor.from_string(h)

    HEAD = "Trebuchet MS"   # sans headings
    BODY = "Georgia"        # serif body

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W, H = 13.333, 7.5
    BLANK = prs.slide_layouts[6]

    def slide(dark=False):
        sl = prs.slides.add_slide(BLANK)
        bg = sl.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb(FOREST if dark else "FFFFFF")
        bg.line.fill.background()
        bg.shadow.inherit = False
        return sl

    def text(sl, s, x, y, w, h, size=16, font=BODY, color=INK, bold=False,
             align=PP_ALIGN.LEFT, italic=False, space_after=6, anchor=None):
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        if anchor is not None:
            tf.vertical_anchor = anchor
        lines = s.split("\n") if isinstance(s, str) else s
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(space_after)
            r = p.add_run()
            r.text = ln
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = rgb(color)
        return tb

    def title(sl, s, dark=False):
        text(sl, s, 0.75, 0.55, W - 1.5, 0.9, size=34, font=HEAD, bold=True,
             color="FFFFFF" if dark else FOREST)

    def bullets(sl, items, x, y, w, size=15, color=INK, gapp=10):
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gapp)
            if isinstance(it, tuple):
                r = p.add_run(); r.text = it[0] + "  "
                r.font.name = HEAD; r.font.size = Pt(size); r.font.bold = True
                r.font.color.rgb = rgb(FOREST if color == INK else color)
                r2 = p.add_run(); r2.text = it[1]
                r2.font.name = BODY; r2.font.size = Pt(size); r2.font.color.rgb = rgb(color)
            else:
                r = p.add_run(); r.text = it
                r.font.name = BODY; r.font.size = Pt(size); r.font.color.rgb = rgb(color)
        return tb

    def card(sl, x, y, w, h, fill=CREAM):
        sh = sl.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
        sh.line.fill.background()
        sh.shadow.inherit = False
        try:
            sh.adjustments[0] = 0.06
        except Exception:
            pass
        return sh

    def stat(sl, value, label, x, y, w=2.85, color=FOREST):
        size = 44 if len(str(value)) <= 6 else 30
        text(sl, value, x, y, w, 1.0, size=size, font=HEAD, bold=True, color=color)
        text(sl, label, x, y + 0.92, w, 0.8, size=12, color=GREY)

    def image(sl, name, x, y, w=None, h=None):
        p = fig(name)
        if p:
            kw = {}
            if w: kw["width"] = Inches(w)
            if h: kw["height"] = Inches(h)
            return sl.shapes.add_picture(p, Inches(x), Inches(y), **kw)
        ph = card(sl, x, y, w or 5.5, h or 3.4, fill="EDEFE9")
        text(sl, f"[{name} appears here\nafter the notebook run]",
             x + 0.3, y + (h or 3.4) / 2 - 0.35, (w or 5.5) - 0.6, 0.8,
             size=12, color=GREY, italic=True, align=PP_ALIGN.CENTER)
        return ph

    def tbl(sl, rows, x, y, w, h, size=12):
        shape = sl.shapes.add_table(len(rows), len(rows[0]),
                                    Inches(x), Inches(y), Inches(w), Inches(h))
        t = shape.table
        for ci in range(len(rows[0])):
            t.columns[ci].width = Emu(int(Inches(w) / len(rows[0])))
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = t.cell(ri, ci)
                cell.text = str(val)
                cell.margin_left = Inches(0.09)
                cell.margin_top = Inches(0.04)
                cell.margin_bottom = Inches(0.04)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(FOREST if ri == 0 else
                                               ("FFFFFF" if ri % 2 else "F1F4EE"))
                p = cell.text_frame.paragraphs[0]
                for r in p.runs:
                    r.font.name = HEAD if ri == 0 else BODY
                    r.font.size = Pt(size)
                    r.font.bold = (ri == 0)
                    r.font.color.rgb = rgb("FFFFFF" if ri == 0 else INK)
        return shape

    # -------------------------------------------------- 1 title
    s1 = slide(dark=True)
    text(s1, TITLE, 0.9, 2.15, 9.6, 1.8, size=48, font=HEAD, bold=True, color="FFFFFF")
    text(s1, SUBTITLE, 0.9, 3.95, 9.6, 0.6, size=19, color=MOSS)
    text(s1, [STUDENT, COURSE + "  |  Capstone project", GITHUB_URL],
         0.9, 5.15, 9.6, 1.4, size=14, color="D8E4D0", space_after=3)
    s1.notes_slide.notes_text_frame.text = (
        "Cassava leaf disease, five classes, real field photos. Two halves: a PyTorch "
        "classifier, then a Claude call that turns the prediction into advice. "
        "Everything runs in a Gradio interface.")

    # -------------------------------------------------- 2 problem
    s2 = slide()
    title(s2, "The problem")
    bullets(s2, [
        ("Why cassava.", "A staple crop across sub-Saharan Africa. Four diseases cause "
                          "most of the yield loss."),
        ("Why it is hard.", "The diseases look alike to a non-expert, and expert diagnosis "
                            "is scarce where the crop matters most."),
        ("Why these images.", "Taken by farmers on phones in the field. Bad light, odd "
                              "angles, cluttered backgrounds."),
        ("Why a label is not enough.", "A five-way softmax is not something a farmer can "
                                       "act on. That is the second half of this project."),
    ], 0.75, 1.85, 6.4, size=15, gapp=16)
    image(s2, "augmented_batch.png", 7.6, 2.1, w=5.0)
    s2.notes_slide.notes_text_frame.text = (
        "Point out that lab-photo plant datasets hit 99 percent without learning much. "
        "Field photos are the harder and more honest version of this task.")

    # -------------------------------------------------- 3 data
    s3 = slide()
    title(s3, "Data")
    bullets(s3, [
        ("iCassava 2019.", "Makerere University AI lab with NaCRRI, Uganda. Labelled by "
                           "plant pathologists."),
        ("9,430 images, 5 classes.", "Fixed split shipped with the data: 5,656 train, "
                                     "1,889 validation, 1,885 test."),
        ("Public download.", "No Kaggle account needed, so the notebook reproduces end to end."),
        ("Not from class or homework.", "Not a built-in torchvision dataset either."),
    ], 0.75, 1.8, 6.3, size=14.5, gapp=14)
    card(s3, 0.75, 5.15, 6.3, 1.25)
    text(s3, "The imbalance is the real difficulty. Mosaic disease has roughly ten times "
             "the images of the healthy class.",
         1.05, 5.43, 5.7, 0.8, size=13.5, italic=True, color=FOREST)
    image(s3, "class_distribution.png", 7.7, 2.2, w=4.9)
    s3.notes_slide.notes_text_frame.text = (
        "The imbalance drives two decisions later: class-weighted loss, and macro F1 "
        "instead of accuracy as the selection metric.")

    # -------------------------------------------------- 4 pipeline
    s4 = slide()
    title(s4, "Preprocessing and augmentation")
    boxes = [
        ("Resize", "224 x 224, ImageNet\nnormalization"),
        ("Geometry", "Random resized crop,\nflips both axes,\nrotation to 20 degrees"),
        ("Colour", "Brightness, contrast,\nsaturation, hue jitter"),
        ("Occlusion", "Random erasing\nat p = 0.25"),
    ]
    for i, (h, b) in enumerate(boxes):
        x = 0.75 + i * 3.15
        card(s4, x, 2.0, 2.85, 2.3)
        text(s4, h, x + 0.28, 2.28, 2.3, 0.45, size=17, font=HEAD, bold=True, color=FOREST)
        text(s4, b, x + 0.28, 2.85, 2.35, 1.3, size=12.5, color=INK, space_after=2)
    card(s4, 0.75, 4.75, 11.85, 1.5)
    text(s4, "Vertical flips are safe here in a way they would not be for text or digits. "
             "A leaf has no canonical orientation.\n"
             "Validation and test do none of this: resize to 256, centre crop to 224, normalize.",
         1.1, 5.05, 11.2, 1.0, size=14, color=INK, space_after=5)
    s4.notes_slide.notes_text_frame.text = (
        "The augmentation is deliberately plain. It is matched to how the photos actually "
        "vary, not chosen to look sophisticated.")

    # -------------------------------------------------- 5 models
    s5 = slide()
    title(s5, "Two models to compare")
    card(s5, 0.75, 1.8, 5.85, 4.3, fill=CREAM)
    text(s5, "SmallCNN", 1.1, 2.1, 5.2, 0.5, size=22, font=HEAD, bold=True, color=FOREST)
    text(s5, "Built from scratch", 1.1, 2.62, 5.2, 0.4, size=13, color=GREY, italic=True)
    bullets(s5, [
        "4 conv blocks, batch norm, dropout 0.4",
        "32 to 256 channels, global average pooling",
        "About 1.2M parameters",
        "The honest baseline: what the data gives you with no outside help",
    ], 1.1, 3.2, 5.15, size=13, gapp=9)
    card(s5, 6.95, 1.8, 5.65, 4.3, fill="E4EDE0")
    text(s5, "ResNet18", 7.3, 2.1, 5.0, 0.5, size=22, font=HEAD, bold=True, color=FOREST)
    text(s5, "ImageNet weights, fine-tuned", 7.3, 2.62, 5.0, 0.4, size=13, color=GREY,
         italic=True)
    bullets(s5, [
        "New 5-way head, whole network trainable",
        "About 11.2M parameters",
        "Under 6,000 training images is too few to learn low-level filters from nothing",
        "Expected to win, and the size of the gap is the result",
    ], 7.3, 3.2, 4.95, size=13, gapp=9)
    text(s5, "Regularization: dropout, weight decay, label smoothing 0.05, augmentation.",
         0.75, 6.35, 11.8, 0.5, size=13, color=GREY, italic=True)
    s5.notes_slide.notes_text_frame.text = (
        "If asked why ResNet18 and not something bigger: 11M parameters trains in minutes on "
        "a free T4, and the point was the comparison, not squeezing the last point of accuracy.")

    # -------------------------------------------------- 6 training
    s6 = slide()
    title(s6, "Training")
    image(s6, "training_curves.png", 0.75, 1.8, w=7.6)
    card(s6, 8.75, 1.8, 3.85, 3.7)
    text(s6, "Choices worth defending", 9.05, 2.05, 3.3, 0.5, size=15, font=HEAD, bold=True,
         color=FOREST)
    bullets(s6, [
        "AdamW with a OneCycle schedule",
        "Class-weighted cross entropy, inverse frequency",
        "Mixed precision, roughly 2x on a T4",
        "Early stopping on validation macro F1, patience 5",
        "Best epoch restored, not the last one",
    ], 9.05, 2.65, 3.25, size=12, gapp=9)
    text(s6, "Accuracy would let the rare classes collapse without ever showing it. "
             "Macro F1 catches that.",
         0.75, 5.85, 7.6, 0.6, size=13.5, italic=True, color=FOREST)
    s6.notes_slide.notes_text_frame.text = (
        "The macro F1 point is the one to make out loud. It is the difference between a "
        "model that looks good and a model that works on the classes you care about.")

    # -------------------------------------------------- 7 sweep
    s7 = slide()
    title(s7, "Hyperparameter tuning")
    text(s7, "Five configurations, six epochs each, ranked on validation macro F1. "
             "The winner was then retrained for up to 25 epochs.",
         0.75, 1.6, 11.8, 0.5, size=14, color=GREY)
    tbl(s7, sweep_rows(), 0.75, 2.35, 7.3, 2.6, size=12)
    card(s7, 8.55, 2.35, 4.05, 1.95)
    text(s7, "Best configuration", 8.9, 2.6, 3.5, 0.4, size=15, font=HEAD, bold=True,
         color=FOREST)
    text(s7, best_cfg_text(), 8.9, 3.1, 3.4, 1.6, size=13.5, color=INK)
    text(s7, "What this does not settle: one seed, six epochs, five points. It picks the "
             "region, not the optimum. Run-to-run variance was never measured.",
         0.75, 5.15, 11.8, 0.9, size=13.5, italic=True, color=FOREST)
    s7.notes_slide.notes_text_frame.text = (
        "Be upfront that the sweep is thin. Saying so is better than being asked.")

    # -------------------------------------------------- 8 results
    s8 = slide()
    title(s8, "Results on the held-out test set")
    stat(s8, pct("resnet18_test_acc"), "ResNet18 accuracy", 0.75, 1.75)
    stat(s8, pct("smallcnn_test_acc"), "SmallCNN accuracy", 3.7, 1.75, color=GREY)
    stat(s8, num("resnet18_test_macro_f1"), "ResNet18 macro F1", 6.65, 1.75)
    stat(s8, gap(), "points gained, transfer vs scratch", 9.6, 1.75, color=MOSS)
    tbl(s8, per_class_rows(), 0.75, 3.55, 7.4, 2.5, size=11.5)
    image(s8, "confusion_test_resnet18.png", 8.65, 3.45, h=2.9)
    s8.notes_slide.notes_text_frame.text = (
        "Lead with the gap, not the accuracy. On a dataset this size, pretrained features "
        "are worth more than any architecture tweak available to me.")

    # -------------------------------------------------- 9 errors
    s9 = slide()
    title(s9, "Where it goes wrong")
    text(s9, "The mistakes worth looking at are the confident ones. Wrong at 40 percent is "
             "at least signalling doubt. Wrong at 95 percent is not.",
         0.75, 1.6, 11.8, 0.5, size=14, color=GREY)
    image(s9, "mistakes.png", 0.75, 2.35, w=11.8)
    s9.notes_slide.notes_text_frame.text = (
        "Replace this note after your run: say which class pairs get confused and whether "
        "any of the errors look like label noise.")

    # -------------------------------------------------- 10 llm + ui
    s10 = slide()
    title(s10, "From prediction to advice")
    bullets(s10, [
        ("The gap.", "A five-way softmax is not actionable. The top three predictions plus "
                     "their probabilities go to Claude through the Anthropic API."),
        ("Why send the probabilities.", "When the classifier is unsure, the response can say "
                                        "so and recommend a human check, instead of turning a "
                                        "coin flip into confident instructions."),
        ("The interface.", "Gradio. Upload a photo, get the class bars and the written "
                           "recommendation side by side. Standalone app.py in the repo too."),
    ], 0.75, 1.85, 6.2, size=14.5, gapp=16)
    image(s10, "ui_screenshot.png", 7.4, 1.9, w=5.2)
    s10.notes_slide.notes_text_frame.text = (
        "Demo live if the share link is up. The five-line API call is the whole integration.")

    # -------------------------------------------------- 11 limits
    s11 = slide(dark=True)
    title(s11, "What this does not prove", dark=True)
    bullets(s11, [
        ("Not generalization.", "Train and test come from the same collection effort. Same "
                                "region, season, and phones. A good test score is not evidence "
                                "it works on another farm."),
        ("Not correct advice.", "No agronomist reviewed a single response. The system writes "
                                "equally fluent advice when the classifier is wrong, which is "
                                "worse than a bare wrong label because it invites action."),
        ("Not a tuned optimum.", "Five configurations, one seed, six epochs. Variance unmeasured."),
    ], 0.75, 2.0, 11.8, size=15.5, color="E8F0E2", gapp=20)
    text(s11, "Fixing the first two needs a field trial and expert review, not more epochs.",
         0.75, 4.75, 11.8, 0.6, size=15, italic=True, color=MOSS)
    s11.notes_slide.notes_text_frame.text = (
        "Say this part plainly. It is the difference between a project that reports a number "
        "and one that understands what the number means.")

    # -------------------------------------------------- 12 next
    s12 = slide()
    title(s12, "Next")
    items = [
        ("Out-of-source test", "Images from a different region or camera, for an honest read "
                               "on generalization."),
        ("Confidence gate", "Below a threshold, route to a human instead of the language model."),
        ("Stronger backbone", "EfficientNet-B0 or ConvNeXt-Tiny, now that the pipeline works."),
        ("Grad-CAM", "Show which part of the leaf drove the call, and catch the model "
                     "latching onto background."),
    ]
    for i, (h, b) in enumerate(items):
        x = 0.75 + (i % 2) * 6.1
        y = 1.95 + (i // 2) * 2.15
        card(s12, x, y, 5.75, 1.85)
        text(s12, h, x + 0.32, y + 0.25, 5.1, 0.45, size=17, font=HEAD, bold=True, color=FOREST)
        text(s12, b, x + 0.32, y + 0.78, 5.15, 0.95, size=12.5, color=INK)
    text(s12, GITHUB_URL, 0.75, 6.5, 11.8, 0.5, size=14, font=HEAD, color=FOREST)
    s12.notes_slide.notes_text_frame.text = "Close on the repo link and take questions."

    prs.save(str(OUT_PPTX))
    print("wrote", OUT_PPTX.name, f"({len(prs.slides._sldIdLst)} slides)")


# ================================================================== zip

def build_zip():
    wanted = ["cassava_capstone.ipynb", "report.pdf", "slides.pptx",
              "README.md", "app.py", "requirements.txt", "make_deliverables.py",
              "results.json"]
    with zipfile.ZipFile(OUT_ZIP, "w", zipfile.ZIP_DEFLATED) as z:
        for f in wanted:
            p = HERE / f
            if p.exists():
                z.write(p, f)
                print("  +", f)
            else:
                print("  ! missing:", f)
        for p in sorted(FIG.glob("*.png")):
            z.write(p, f"figures/{p.name}")
    print("wrote", OUT_ZIP.name, f"({OUT_ZIP.stat().st_size/1e6:.1f} MB)")


if __name__ == "__main__":
    build_pdf()
    build_pptx()
    build_zip()
    if GITHUB_URL == "GITHUB_LINK_HERE":
        print("\n>>> GITHUB_URL is still a placeholder. Set it at the top of this file "
              "and run again before submitting.")
